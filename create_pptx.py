#!/usr/bin/env python3
"""
Generate Maersk SF Task Force Report-out PowerPoint
Matches the Bosch "Project Status | General" template from the screenshot.
Populated with KPI data from CW01, CW03, CW05 analysis.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors ──────────────────────────────────────────────────────────────────
BOSCH_RED = RGBColor(0xE2, 0x00, 0x15)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN = RGBColor(0x00, 0xB0, 0x50)
YELLOW = RGBColor(0xFF, 0xC0, 0x00)
RED = RGBColor(0xFF, 0x00, 0x00)
HEADER_BG = RGBColor(0x00, 0x2B, 0x5C)  # Maersk dark blue
MAERSK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xE4, 0xF0)
TABLE_HEADER_BG = RGBColor(0x00, 0x2B, 0x5C)

# ─── KPI Data (from pdca_kpi_extract.py output) ─────────────────────────────
KPI_DATA = {
    "CW01": {
        "crit_comp": 82.4, "crit_comp_sc3": 77.4, "crit_comp_sc4": 91.1,
        "crit_time": 52.3, "crit_time_sc3": 58.3, "crit_time_sc4": 41.6,
        "eta_2p": 58.3, "eta_2d": 34.7,
        "plausibility": 75.1,
        "all_comp": 72.7, "all_comp_sc3": 68.5, "all_comp_sc4": 77.3,
        "all_time": 45.8, "all_time_sc3": 48.4, "all_time_sc4": 42.9,
        "ref_comp": 0.0,
    },
    "CW03": {
        "crit_comp": 80.0, "crit_comp_sc3": 75.5, "crit_comp_sc4": 82.2,
        "crit_time": 48.9, "crit_time_sc3": 54.6, "crit_time_sc4": 46.1,
        "eta_2p": 42.9, "eta_2d": 42.6,
        "plausibility": 66.4,
        "all_comp": 65.4, "all_comp_sc3": 69.1, "all_comp_sc4": 64.2,
        "all_time": 37.9, "all_time_sc3": 48.9, "all_time_sc4": 34.2,
        "ref_comp": 0.0,
    },
    "CW05": {
        "crit_comp": 83.9, "crit_comp_sc3": 80.7, "crit_comp_sc4": 85.9,
        "crit_time": 45.4, "crit_time_sc3": 60.7, "crit_time_sc4": 35.6,
        "eta_2p": 54.4, "eta_2d": 47.2,
        "plausibility": 67.1,
        "all_comp": 74.7, "all_comp_sc3": 76.1, "all_comp_sc4": 74.2,
        "all_time": 45.1, "all_time_sc3": 57.7, "all_time_sc4": 39.6,
        "ref_comp": 0.0,
    },
}

TARGETS = {
    "crit_comp": 95, "crit_time": 70, "eta_2p": 70, "eta_2d": 70,
    "plausibility": 90, "all_comp": 95, "all_time": 70, "ref_comp": 95,
}


def get_status_color(value, target):
    """Return RAG color based on value vs target."""
    if value >= target:
        return GREEN
    elif value >= target * 0.8:
        return YELLOW
    else:
        return RED


def get_status_letter(value, target):
    if value >= target:
        return "G"
    elif value >= target * 0.8:
        return "Y"
    else:
        return "R"


def get_trend(kpi_key):
    """Return trend arrow based on CW03 -> CW05."""
    v3 = KPI_DATA["CW03"][kpi_key]
    v5 = KPI_DATA["CW05"][kpi_key]
    if v5 > v3 + 1:
        return "▲"
    elif v5 < v3 - 1:
        return "▼"
    else:
        return "►"


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=10,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_circle(slide, left, top, size, fill_color, text="", font_size=8):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = Inches(13.333)
    slide_h = Inches(7.5)

    # =========================================================================
    # SLIDE 1: Report-out cover (matching the screenshot template)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # --- Top red bar ---
    add_shape(slide, Inches(0), Inches(0), slide_w, Inches(0.08), BOSCH_RED)

    # --- Header section ---
    add_shape(slide, Inches(0), Inches(0.08), slide_w, Inches(0.8), MAERSK_BLUE)
    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(8), Inches(0.65),
                 "Maersk SF Task Force | Report-out", 24, True, WHITE, font_name="Arial")

    # --- Sub-header row ---
    sub_y = Inches(0.88)
    add_shape(slide, Inches(0), sub_y, slide_w, Inches(0.55), LIGHT_BLUE_BG)

    add_text_box(slide, Inches(0.4), sub_y + Pt(4), Inches(2), Inches(0.45),
                 "Work Stream:", 11, True, DARK_GRAY)
    add_text_box(slide, Inches(2.0), sub_y + Pt(4), Inches(3), Inches(0.45),
                 "EDI - Visibility", 11, False, DARK_GRAY)

    add_text_box(slide, Inches(5.5), sub_y + Pt(4), Inches(2), Inches(0.45),
                 "Responsible:", 11, True, DARK_GRAY)
    add_text_box(slide, Inches(7.2), sub_y + Pt(4), Inches(2), Inches(0.45),
                 "Murali Rajamani", 11, False, DARK_GRAY)

    # Overall status
    add_text_box(slide, Inches(10.0), sub_y + Pt(4), Inches(1.2), Inches(0.45),
                 "Status:", 11, True, DARK_GRAY)
    # Determine overall status: critical completeness at 83.9% (target 95) = Yellow
    overall_status = get_status_letter(KPI_DATA["CW05"]["crit_comp"], TARGETS["crit_comp"])
    overall_color = get_status_color(KPI_DATA["CW05"]["crit_comp"], TARGETS["crit_comp"])
    add_circle(slide, Inches(10.9), sub_y + Pt(4), Inches(0.35), overall_color, overall_status, 12)

    add_text_box(slide, Inches(11.8), sub_y + Pt(4), Inches(1.2), Inches(0.45),
                 "Date: 26.02.2026", 11, False, DARK_GRAY)

    # --- Important milestones (traffic lights) ---
    mile_y = Inches(1.55)
    add_shape(slide, Inches(0), mile_y, slide_w, Inches(0.55), LIGHT_GRAY)
    add_text_box(slide, Inches(0.4), mile_y + Pt(2), Inches(2), Inches(0.45),
                 "Important milestones", 10, True, DARK_GRAY)

    milestone_labels = [
        ("Completeness", KPI_DATA["CW05"]["crit_comp"], TARGETS["crit_comp"]),
        ("Timeliness", KPI_DATA["CW05"]["crit_time"], TARGETS["crit_time"]),
        ("ETA to Port", KPI_DATA["CW05"]["eta_2p"], TARGETS["eta_2p"]),
        ("ETA to Door", KPI_DATA["CW05"]["eta_2d"], TARGETS["eta_2d"]),
        ("Plausibility", KPI_DATA["CW05"]["plausibility"], TARGETS["plausibility"]),
        ("Ref. Comp.", KPI_DATA["CW05"]["ref_comp"], TARGETS["ref_comp"]),
    ]

    x_start = Inches(3.0)
    for i, (label, val, tgt) in enumerate(milestone_labels):
        x_pos = x_start + Inches(i * 1.7)
        color = get_status_color(val, tgt)
        letter = get_status_letter(val, tgt)
        add_circle(slide, x_pos, mile_y + Pt(5), Inches(0.3), color, letter, 10)
        add_text_box(slide, x_pos + Inches(0.35), mile_y + Pt(3), Inches(1.3), Inches(0.45),
                     label, 9, False, DARK_GRAY)

    # =========================================================================
    # LEFT COLUMN: Done + Planned
    # =========================================================================
    col_left = Inches(0.3)
    col_width = Inches(6.3)

    # --- Done (last period) ---
    done_y = Inches(2.3)
    add_shape(slide, col_left, done_y, col_width, Inches(0.35), MAERSK_BLUE)
    add_text_box(slide, col_left + Pt(4), done_y + Pt(2), col_width, Inches(0.3),
                 "Done (last period) — CW05 Results", 11, True, WHITE)

    done_items = [
        f"Critical Milestone Completeness improved to 83.9% (CW03: 80.0%, +3.9pp)",
        f"SC3 Critical Completeness: 80.7% (+5.2pp vs CW03)",
        f"SC4 Critical Completeness: 85.9% (+3.7pp vs CW03)",
        f"All Milestones Completeness improved to 74.7% (CW03: 65.4%, +9.3pp)",
        f"ETA Accuracy to Door improved to 47.2% (CW03: 42.6%, +4.6pp)",
        f"SC3 Critical Timeliness improved to 60.7% (CW03: 54.6%, +6.1pp)",
        f"Plausibility stabilized at 67.1% (CW03: 66.4%)",
    ]

    done_text_y = done_y + Inches(0.4)
    for i, item in enumerate(done_items):
        add_text_box(slide, col_left + Pt(8), done_text_y + Inches(i * 0.22), col_width - Pt(16), Inches(0.22),
                     f"• {item}", 9, False, DARK_GRAY)

    # --- Planned activities for next period ---
    plan_y = done_text_y + Inches(len(done_items) * 0.22) + Inches(0.12)
    add_shape(slide, col_left, plan_y, col_width, Inches(0.35), MAERSK_BLUE)
    add_text_box(slide, col_left + Pt(4), plan_y + Pt(2), col_width, Inches(0.3),
                 "Planned activities for next period", 11, True, WHITE)

    plan_items = [
        "Drive SC4 Timeliness improvement (currently 35.6% — lowest KPI)",
        "Investigate SC4 S00 Shipment Created completeness (55.2%)",
        "Improve S46 Docs Rcvd from Shipper (67.7% completeness)",
        "Address SC4 S05 In-delivery gap (17.5% completeness)",
        "Target Plausibility improvement toward 90% (currently 67.1%)",
        "Root-cause S54 Full Container Discharge failure (7.4% SC3)",
        "Establish Reference Completeness baseline (currently 0%)",
    ]

    plan_text_y = plan_y + Inches(0.4)
    for i, item in enumerate(plan_items):
        add_text_box(slide, col_left + Pt(8), plan_text_y + Inches(i * 0.22), col_width - Pt(16), Inches(0.22),
                     f"• {item}", 9, False, DARK_GRAY)

    # =========================================================================
    # RIGHT COLUMN: Critical issues
    # =========================================================================
    col_right = Inches(6.9)
    col_right_w = Inches(6.1)

    issue_y = Inches(2.3)
    add_shape(slide, col_right, issue_y, col_right_w, Inches(0.35), BOSCH_RED)
    add_text_box(slide, col_right + Pt(4), issue_y + Pt(2), col_right_w, Inches(0.3),
                 "Critical issues / Problems", 11, True, WHITE)

    # Table headers
    th_y = issue_y + Inches(0.38)
    issue_col_w = Inches(2.6)
    sol_col_w = Inches(2.0)
    resp_col_w = Inches(0.9)
    date_col_w = Inches(0.6)

    add_shape(slide, col_right, th_y, col_right_w, Inches(0.28), LIGHT_GRAY)
    add_text_box(slide, col_right + Pt(2), th_y + Pt(1), issue_col_w, Inches(0.25),
                 "Issue", 9, True, DARK_GRAY)
    add_text_box(slide, col_right + issue_col_w, th_y + Pt(1), sol_col_w, Inches(0.25),
                 "Solution", 9, True, DARK_GRAY)
    add_text_box(slide, col_right + issue_col_w + sol_col_w, th_y + Pt(1), resp_col_w, Inches(0.25),
                 "Resp.", 9, True, DARK_GRAY)
    add_text_box(slide, col_right + issue_col_w + sol_col_w + resp_col_w, th_y + Pt(1), date_col_w, Inches(0.25),
                 "Until", 9, True, DARK_GRAY)

    issues = [
        ("SC4 Critical Timeliness at 35.6%\n(target: 70%)", "Analyze late milestone\ntransmission root causes", "Team", "CW07"),
        ("SC4 S05 In-delivery only 17.5%", "Review carrier EDI\nintegration gaps", "Team", "CW07"),
        ("Plausibility at 67.1%\n(target: 90%)", "Identify out-of-sequence\nmilestone patterns", "Team", "CW08"),
        ("S54 Container Discharge 7.4%", "Escalate terminal EDI\nfeed issues", "Team", "CW07"),
        ("Reference Completeness 0%", "Define REF field mapping\nrequirements with Bosch", "Team", "CW08"),
        ("ETA to Door accuracy 47.2%", "Calibrate estimated\ndelivery models", "Team", "CW09"),
    ]

    row_y = th_y + Inches(0.30)
    row_h = Inches(0.42)
    for idx, (issue, solution, resp, until) in enumerate(issues):
        bg = WHITE if idx % 2 == 0 else LIGHT_BLUE_BG
        add_shape(slide, col_right, row_y, col_right_w, row_h, bg, LIGHT_GRAY, Pt(0.5))
        add_text_box(slide, col_right + Pt(2), row_y + Pt(1), issue_col_w, row_h,
                     issue, 8, False, DARK_GRAY)
        add_text_box(slide, col_right + issue_col_w, row_y + Pt(1), sol_col_w, row_h,
                     solution, 8, False, DARK_GRAY)
        add_text_box(slide, col_right + issue_col_w + sol_col_w, row_y + Pt(1), resp_col_w, row_h,
                     resp, 8, False, DARK_GRAY, PP_ALIGN.CENTER)
        add_text_box(slide, col_right + issue_col_w + sol_col_w + resp_col_w, row_y + Pt(1), date_col_w, row_h,
                     until, 8, False, DARK_GRAY, PP_ALIGN.CENTER)
        row_y += row_h

    # --- Footer ---
    add_shape(slide, Inches(0), slide_h - Inches(0.35), slide_w, Inches(0.35), MAERSK_BLUE)
    add_text_box(slide, Inches(0.4), slide_h - Inches(0.33), Inches(5), Inches(0.3),
                 "Bosch | Project Status | General", 8, False, WHITE)
    add_text_box(slide, Inches(10), slide_h - Inches(0.33), Inches(3), Inches(0.3),
                 "26.02.2026 | Confidential", 8, False, WHITE, PP_ALIGN.RIGHT)

    # =========================================================================
    # SLIDE 2: KPI Trend Dashboard
    # =========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Top bar
    add_shape(slide2, Inches(0), Inches(0), slide_w, Inches(0.08), BOSCH_RED)
    add_shape(slide2, Inches(0), Inches(0.08), slide_w, Inches(0.7), MAERSK_BLUE)
    add_text_box(slide2, Inches(0.4), Inches(0.15), Inches(10), Inches(0.55),
                 "Bosch Milestone KPI Trend — CW01 / CW03 / CW05", 22, True, WHITE)

    # --- KPI Summary Table ---
    table_x = Inches(0.4)
    table_y = Inches(1.1)
    cols = 6
    rows = 10
    col_widths = [Inches(3.2), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.0)]
    row_height = Inches(0.42)

    # Headers
    headers = ["KPI", "Target", "CW01", "CW03", "CW05", "Trend"]
    for c, (hdr, cw) in enumerate(zip(headers, col_widths)):
        x = table_x + sum(w for w in [v for v in col_widths[:c]])
        add_shape(slide2, x, table_y, cw, Inches(0.38), MAERSK_BLUE)
        add_text_box(slide2, x + Pt(4), table_y + Pt(2), cw - Pt(8), Inches(0.34),
                     hdr, 10, True, WHITE, PP_ALIGN.CENTER)

    kpi_rows = [
        ("P1: Critical Completeness", "≥95%", "crit_comp"),
        ("   SC3 Critical Comp.", "", "crit_comp_sc3"),
        ("   SC4 Critical Comp.", "", "crit_comp_sc4"),
        ("P2: Critical Timeliness", "≥70%", "crit_time"),
        ("   SC3 Critical Time.", "", "crit_time_sc3"),
        ("   SC4 Critical Time.", "", "crit_time_sc4"),
        ("P3: ETA Accuracy to Port", "≥70%", "eta_2p"),
        ("P4: ETA Accuracy to Door", "≥70%", "eta_2d"),
        ("P7: Plausibility", "≥90%", "plausibility"),
        ("P8: All MS Completeness", "≥95%", "all_comp"),
        ("P9: All MS Timeliness", "≥70%", "all_time"),
    ]

    for r, (label, target, key) in enumerate(kpi_rows):
        ry = table_y + Inches(0.40) + r * row_height
        bg = WHITE if r % 2 == 0 else LIGHT_BLUE_BG
        is_sub = label.startswith("   ")

        for c, cw in enumerate(col_widths):
            x = table_x + sum(v for v in col_widths[:c])
            add_shape(slide2, x, ry, cw, row_height, bg, LIGHT_GRAY, Pt(0.5))

        # Label
        x0 = table_x
        fsize = 9 if is_sub else 10
        fbold = not is_sub
        add_text_box(slide2, x0 + Pt(4), ry + Pt(3), col_widths[0] - Pt(8), row_height - Pt(6),
                     label, fsize, fbold, DARK_GRAY)

        # Target
        x1 = table_x + col_widths[0]
        add_text_box(slide2, x1 + Pt(4), ry + Pt(3), col_widths[1] - Pt(8), row_height - Pt(6),
                     target, 9, False, DARK_GRAY, PP_ALIGN.CENTER)

        # CW values
        for ci, cw_name in enumerate(["CW01", "CW03", "CW05"]):
            val = KPI_DATA[cw_name].get(key, None)
            if val is not None:
                x_col = table_x + sum(v for v in col_widths[:ci + 2])
                val_text = f"{val:.1f}%"
                # Color code
                parent_key = key.replace("_sc3", "").replace("_sc4", "")
                tgt_val = TARGETS.get(parent_key, TARGETS.get(key, 95))
                val_color = get_status_color(val, tgt_val)
                add_text_box(slide2, x_col + Pt(4), ry + Pt(3), col_widths[ci + 2] - Pt(8), row_height - Pt(6),
                             val_text, 10, True, val_color, PP_ALIGN.CENTER)

        # Trend arrow
        x_trend = table_x + sum(v for v in col_widths[:5])
        trend = get_trend(key) if key in KPI_DATA["CW03"] else ""
        trend_color = GREEN if "▲" in trend else (RED if "▼" in trend else YELLOW)
        add_text_box(slide2, x_trend + Pt(4), ry + Pt(3), col_widths[5] - Pt(8), row_height - Pt(6),
                     trend, 14, True, trend_color, PP_ALIGN.CENTER)

    # --- Key Takeaways ---
    takeaway_y = table_y + Inches(0.40) + len(kpi_rows) * row_height + Inches(0.3)
    add_shape(slide2, Inches(0.4), takeaway_y, Inches(12.5), Inches(0.35), MAERSK_BLUE)
    add_text_box(slide2, Inches(0.6), takeaway_y + Pt(2), Inches(12), Inches(0.3),
                 "Key Takeaways", 11, True, WHITE)

    takeaways = [
        "Completeness improving: Critical +3.9pp (CW03→CW05), All Milestones +9.3pp — strongest improvement area",
        "Timeliness remains the biggest challenge: SC4 Critical at 35.6%, overall trend declining from CW01 (52.3%→45.4%)",
        "ETA accuracy mixed: Port accuracy dropped from CW01 (58.3%→54.4%), Door accuracy improved (34.7%→47.2%)",
        "SC3 consistently outperforms SC4 across all metrics — SC4 operational gaps driving overall underperformance",
    ]

    for i, item in enumerate(takeaways):
        add_text_box(slide2, Inches(0.6), takeaway_y + Inches(0.4) + Inches(i * 0.25),
                     Inches(12), Inches(0.25), f"• {item}", 9, False, DARK_GRAY)

    # Footer
    add_shape(slide2, Inches(0), slide_h - Inches(0.35), slide_w, Inches(0.35), MAERSK_BLUE)
    add_text_box(slide2, Inches(0.4), slide_h - Inches(0.33), Inches(5), Inches(0.3),
                 "Bosch | KPI Dashboard", 8, False, WHITE)
    add_text_box(slide2, Inches(10), slide_h - Inches(0.33), Inches(3), Inches(0.3),
                 "26.02.2026 | Confidential", 8, False, WHITE, PP_ALIGN.RIGHT)

    # =========================================================================
    # SLIDE 3: Detailed Milestone Breakdown
    # =========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_shape(slide3, Inches(0), Inches(0), slide_w, Inches(0.08), BOSCH_RED)
    add_shape(slide3, Inches(0), Inches(0.08), slide_w, Inches(0.7), MAERSK_BLUE)
    add_text_box(slide3, Inches(0.4), Inches(0.15), Inches(10), Inches(0.55),
                 "CW05 Critical Milestone Detail — SC3 vs SC4", 22, True, WHITE)

    # SC3 Key milestones
    sc3_y = Inches(1.1)
    add_shape(slide3, Inches(0.4), sc3_y, Inches(6.1), Inches(0.35), MAERSK_BLUE)
    add_text_box(slide3, Inches(0.6), sc3_y + Pt(2), Inches(5.5), Inches(0.3),
                 "SC3 (FCL) Key Milestones — CW05", 11, True, WHITE)

    sc3_milestones = [
        ("S60 Pre-Booking", 92.3, 92.3),
        ("S02 Collected", 75.7, 34.2),
        ("S04 Departed", 91.2, 89.1),
        ("S07 Arrived", 93.3, 93.0),
        ("S31 Delivered", 56.6, 12.8),
        ("S52 Empty Pickup", 28.9, 24.4),
        ("S53 Full Loaded", 68.5, 63.0),
        ("S54 Full Discharged", 7.4, 5.6),
        ("S55 Empty Return", 76.9, 76.9),
    ]

    sc3_hdr_y = sc3_y + Inches(0.38)
    sc3_hdrs = ["Milestone", "Comp%", "Time%"]
    sc3_cw = [Inches(2.2), Inches(1.8), Inches(1.8)]
    for c, (h, w) in enumerate(zip(sc3_hdrs, sc3_cw)):
        x = Inches(0.4) + sum(v for v in sc3_cw[:c])
        add_shape(slide3, x, sc3_hdr_y, w, Inches(0.3), LIGHT_GRAY)
        add_text_box(slide3, x + Pt(4), sc3_hdr_y + Pt(2), w - Pt(8), Inches(0.26),
                     h, 9, True, DARK_GRAY, PP_ALIGN.CENTER)

    for r, (ms, comp, time) in enumerate(sc3_milestones):
        ry = sc3_hdr_y + Inches(0.32) + r * Inches(0.3)
        bg = WHITE if r % 2 == 0 else LIGHT_BLUE_BG
        for c, w in enumerate(sc3_cw):
            x = Inches(0.4) + sum(v for v in sc3_cw[:c])
            add_shape(slide3, x, ry, w, Inches(0.3), bg, LIGHT_GRAY, Pt(0.5))

        add_text_box(slide3, Inches(0.4) + Pt(4), ry + Pt(2), sc3_cw[0] - Pt(8), Inches(0.26),
                     ms, 9, False, DARK_GRAY)
        comp_color = get_status_color(comp, 95)
        add_text_box(slide3, Inches(0.4) + sc3_cw[0] + Pt(4), ry + Pt(2), sc3_cw[1] - Pt(8), Inches(0.26),
                     f"{comp:.1f}%", 9, True, comp_color, PP_ALIGN.CENTER)
        time_color = get_status_color(time, 70)
        add_text_box(slide3, Inches(0.4) + sc3_cw[0] + sc3_cw[1] + Pt(4), ry + Pt(2), sc3_cw[2] - Pt(8), Inches(0.26),
                     f"{time:.1f}%", 9, True, time_color, PP_ALIGN.CENTER)

    # SC4 Key milestones
    sc4_x = Inches(7.0)
    add_shape(slide3, sc4_x, sc3_y, Inches(6.0), Inches(0.35), MAERSK_BLUE)
    add_text_box(slide3, sc4_x + Inches(0.2), sc3_y + Pt(2), Inches(5.5), Inches(0.3),
                 "SC4 (LCL) Key Milestones — CW05", 11, True, WHITE)

    sc4_milestones = [
        ("S00 Shipment Created", 55.2, 54.6),
        ("S02 Collected", 88.9, 32.9),
        ("S04 Departed", 89.4, 49.1),
        ("S07 Arrived", 81.8, 43.8),
        ("S16 Booked w/ Carrier", 76.9, 61.1),
        ("S17 Tendered Carrier", 59.5, 32.8),
        ("S31 Delivered", 84.0, 10.5),
        ("S45 Handover Broker", 47.7, 47.7),
        ("S46 Docs Rcvd", 67.7, 31.3),
    ]

    sc4_cw = [Inches(2.2), Inches(1.7), Inches(1.7)]
    sc4_hdr_y = sc3_y + Inches(0.38)
    for c, (h, w) in enumerate(zip(sc3_hdrs, sc4_cw)):
        x = sc4_x + sum(v for v in sc4_cw[:c])
        add_shape(slide3, x, sc4_hdr_y, w, Inches(0.3), LIGHT_GRAY)
        add_text_box(slide3, x + Pt(4), sc4_hdr_y + Pt(2), w - Pt(8), Inches(0.26),
                     h, 9, True, DARK_GRAY, PP_ALIGN.CENTER)

    for r, (ms, comp, time) in enumerate(sc4_milestones):
        ry = sc4_hdr_y + Inches(0.32) + r * Inches(0.3)
        bg = WHITE if r % 2 == 0 else LIGHT_BLUE_BG
        for c, w in enumerate(sc4_cw):
            x = sc4_x + sum(v for v in sc4_cw[:c])
            add_shape(slide3, x, ry, w, Inches(0.3), bg, LIGHT_GRAY, Pt(0.5))

        add_text_box(slide3, sc4_x + Pt(4), ry + Pt(2), sc4_cw[0] - Pt(8), Inches(0.26),
                     ms, 9, False, DARK_GRAY)
        comp_color = get_status_color(comp, 95)
        add_text_box(slide3, sc4_x + sc4_cw[0] + Pt(4), ry + Pt(2), sc4_cw[1] - Pt(8), Inches(0.26),
                     f"{comp:.1f}%", 9, True, comp_color, PP_ALIGN.CENTER)
        time_color = get_status_color(time, 70)
        add_text_box(slide3, sc4_x + sc4_cw[0] + sc4_cw[1] + Pt(4), ry + Pt(2), sc4_cw[2] - Pt(8), Inches(0.26),
                     f"{time:.1f}%", 9, True, time_color, PP_ALIGN.CENTER)

    # Bottom summary
    sum_y = sc4_hdr_y + Inches(0.32) + len(sc4_milestones) * Inches(0.3) + Inches(0.3)
    add_shape(slide3, Inches(0.4), sum_y, Inches(12.5), Inches(0.35), MAERSK_BLUE)
    add_text_box(slide3, Inches(0.6), sum_y + Pt(2), Inches(12), Inches(0.3),
                 "Focus Areas", 11, True, WHITE)

    focus_items = [
        "RED — S54 Full Container Discharge (SC3): 7.4% completeness — terminal EDI feed likely broken/missing",
        "RED — S05 In-delivery (SC4): 17.5% completeness — last-mile tracking gap for LCL shipments",
        "RED — S52 Empty Container Pickup (SC3): 28.9% — depot EDI not flowing consistently",
        "RED — S31 Delivered Timeliness: SC3 12.8%, SC4 10.5% — delivery confirmation significantly delayed",
        "YELLOW — S00 Shipment Created (SC4): 55.2% — shipment creation not captured for ~half of shipments",
    ]

    for i, item in enumerate(focus_items):
        add_text_box(slide3, Inches(0.6), sum_y + Inches(0.4) + Inches(i * 0.23),
                     Inches(12), Inches(0.23), f"• {item}", 9, False, DARK_GRAY)

    # Footer
    add_shape(slide3, Inches(0), slide_h - Inches(0.35), slide_w, Inches(0.35), MAERSK_BLUE)
    add_text_box(slide3, Inches(0.4), slide_h - Inches(0.33), Inches(5), Inches(0.3),
                 "Bosch | Milestone Detail", 8, False, WHITE)
    add_text_box(slide3, Inches(10), slide_h - Inches(0.33), Inches(3), Inches(0.3),
                 "26.02.2026 | Confidential", 8, False, WHITE, PP_ALIGN.RIGHT)

    # ─── Save ────────────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__),
                               "Maersk_SF_TaskForce_Reportout_CW05.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Slide 1: Report-out (matching template)")
    print(f"  Slide 2: KPI Trend Dashboard (CW01/CW03/CW05)")
    print(f"  Slide 3: Detailed Milestone Breakdown (SC3 vs SC4)")


if __name__ == "__main__":
    create_presentation()
